from pathlib import Path
import re
import tempfile
from typing import Dict, List

from services.analysis_service import AnalysisService
from services.document_service import DocumentService
from services.llm_service import LLMService

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.util import Cm, Inches, Pt
except ImportError:  # pragma: no cover - handled at runtime
    Presentation = None

from datetime import datetime


BASE_DIR = Path(__file__).resolve().parent.parent
REPORT_DIR = BASE_DIR / "data" / "reports"
FALLBACK_REPORT_DIR = BASE_DIR / "reports"
TEMP_REPORT_DIR = Path(tempfile.gettempdir()) / "tax_monitor_reports"


class ReportService:
    def __init__(self):
        self.document_service = DocumentService()
        self.analysis_service = AnalysisService()
        self.llm_service = LLMService()

    async def generate_report(
        self,
        doc_id: str,
        output_format: str = "obsidian",
        provider: str = "ollama",
        model_name: str = "qwen3:8b",
        target_language: str = "zh",
        user_prompt: str = None
    ) -> Dict:
        document = self.document_service.get_document(doc_id)
        if not document:
            raise ValueError("Document not found")

        try:
            analysis = await self.analysis_service.analyze_document(
                doc_id=doc_id,
                mode="translate_first",
                target_language=target_language,
                use_llm=True,
                provider=provider,
                user_prompt=user_prompt,
                model_name=model_name
            )
        except Exception:
            analysis = await self.analysis_service.analyze_document(
                doc_id=doc_id,
                mode="translate_first",
                target_language=target_language,
                use_llm=False,
                provider=provider,
                user_prompt=user_prompt,
                model_name=model_name
            )

        title = analysis["title"]
        slide_outline = self._build_slide_outline(
            analysis=analysis,
            provider=provider,
            model_name=model_name
        )
        file_path = None

        if output_format == "pptx":
            sources_slide = self._build_sources_slide(
                document=document,
                analysis=analysis,
                provider=provider,
                model_name=model_name,
                target_language=target_language,
                user_prompt=user_prompt,
            )
            slide_outline_with_sources = list(slide_outline) + [sources_slide]
            file_path = self._build_pptx_file(title, slide_outline_with_sources)
            content = f"PPTX report generated: {file_path}"
        elif output_format == "markdown":
            content = self._build_markdown_report(
                title=title,
                document=document,
                analysis=analysis,
                slide_outline=slide_outline,
                provider=provider,
                model_name=model_name,
                target_language=target_language,
                user_prompt=user_prompt,
            )
            file_path = self._write_markdown_file(title, content)
        elif output_format == "slides":
            content = self._build_slide_markdown(title, slide_outline)
        else:
            content = self._build_obsidian_note(title, analysis, slide_outline)

        return {
            "doc_id": doc_id,
            "output_format": output_format,
            "title": title,
            "content": content,
            "slide_outline": slide_outline,
            "file_path": file_path,
            "metadata": {
                "risk_level": analysis["risk_level"],
                "target_language": target_language,
                "provider": provider,
                "model_name": model_name
            }
        }

    def _build_slide_outline(self, analysis: Dict, provider: str, model_name: str) -> List[Dict]:
        fallback = [
            {"title": "監測摘要", "bullets": [analysis["summary"]]},
            {"title": "風險判斷", "bullets": [f"風險等級：{analysis['risk_level']}"]},
            {"title": "關鍵證據", "bullets": analysis.get("evidence", [])[:3]},
            {"title": "建議行動", "bullets": self._build_action_items(analysis)}
        ]

        prompt = f"""
你是一位稅務顧問，請將以下分析內容整理成幾頁簡報大綱。要有地區分別、重點摘要、風險判斷與建議行動。
策略要盡量精準而不籠統。
每頁都要有 title 與 bullets。
只輸出 JSON：
{{
  "slides": [
    {{"title": "頁1", "bullets": ["重點1", "重點2"]}}
  ]
}}

分析內容：
標題：{analysis['title']}
風險等級：{analysis['risk_level']}
摘要：{analysis['summary']}
證據：{analysis.get('evidence', [])}
備註：{analysis.get('notes', [])}
"""
        data = self.llm_service.generate_json(
            prompt=prompt,
            schema_hint={"slides": fallback},
            provider=provider,
            model_name=model_name
        )
        slides = data.get("slides", fallback)
        normalized = []
        for slide in slides[:4]:
            title = slide.get("title") if isinstance(slide, dict) else None
            bullets = slide.get("bullets") if isinstance(slide, dict) else None
            normalized.append({
                "title": title or "未命名頁",
                "bullets": bullets or []
            })
        return normalized or fallback

    def _build_markdown_report(
        self,
        title: str,
        document: Dict,
        analysis: Dict,
        slide_outline: List[Dict],
        provider: str,
        model_name: str,
        target_language: str,
        user_prompt: str = None,
    ) -> str:
        lines = [f"# {title}", ""]
        lines.append("## Executive summary")
        lines.append(analysis.get("summary") or "")
        lines.append("")
        lines.append(f"**Risk level:** {analysis.get('risk_level', 'n/a')}")
        risk_tags = ", ".join(analysis.get("risk_tags", []) or [])
        if risk_tags:
            lines.append(f"**Risk tags:** {risk_tags}")
        lines.append("")
        lines.append("## Auto-extracted keywords")
        keywords = analysis.get("auto_keywords", []) or []
        if keywords:
            lines.append(", ".join(keywords))
        else:
            lines.append("_(none)_")
        lines.append("")
        lines.append("## Key evidence")
        evidence = analysis.get("evidence", []) or []
        if evidence:
            lines.extend(f"- {item}" for item in evidence)
        else:
            lines.append("_(none)_")
        lines.append("")
        if analysis.get("notes"):
            lines.append("## Notes")
            lines.extend(f"- {note}" for note in analysis["notes"])
            lines.append("")
        lines.append("## Slide outline")
        for index, slide in enumerate(slide_outline, start=1):
            lines.append(f"### Slide {index}: {slide.get('title', 'Untitled')}")
            lines.extend(f"- {bullet}" for bullet in slide.get("bullets", []))
            lines.append("")
        lines.append("## Sources & metadata")
        lines.append(f"- **Source title:** {document.get('title', 'n/a')}")
        if document.get("url"):
            lines.append(f"- **URL:** <{document['url']}>")
        if document.get("file_name"):
            lines.append(f"- **File name:** {document['file_name']}")
        lines.append(
            f"- **Source type / channel:** {document.get('source_type', 'n/a')} via {document.get('source_name', 'n/a')}"
        )
        lines.append(f"- **Original language:** {document.get('language', 'n/a')}")
        if document.get("country") or document.get("industry"):
            lines.append(f"- **Country / industry:** {document.get('country') or 'n/a'} · {document.get('industry') or 'n/a'}")
        if document.get("published_date"):
            lines.append(f"- **Published:** {document['published_date']}")
        if document.get("created_at"):
            lines.append(f"- **Ingested at:** {document['created_at']}")
        lines.append(f"- **Generated by:** provider=`{provider}`, model=`{model_name}`, target_language=`{target_language}`")
        lines.append(f"- **Report timestamp:** {datetime.now().isoformat(timespec='seconds')}")
        if user_prompt:
            short_prompt = user_prompt if len(user_prompt) <= 240 else user_prompt[:240] + "…"
            lines.append(f"- **Research intent:** {short_prompt}")
        return "\n".join(lines).strip() + "\n"

    def _write_markdown_file(self, title: str, content: str) -> str:
        report_dir = self._get_report_dir()
        file_path = report_dir / f"{self._sanitize_file_name(title)}.md"
        file_path.write_text(content, encoding="utf-8")
        return str(file_path)

    def _build_obsidian_note(self, title: str, analysis: Dict, slide_outline: List[Dict]) -> str:
        frontmatter = [
            "---",
            f"title: {title}",
            "tags:",
            "  - tax-monitor",
            f"  - risk-{analysis['risk_level'].lower()}",
            "---",
            ""
        ]
        body = [
            f"# {title}",
            "",
            "## 監測摘要",
            analysis["summary"],
            "",
            "## 自動關鍵字",
            ", ".join(analysis.get("auto_keywords", [])),
            "",
            "## 風險等級",
            analysis["risk_level"],
            "",
            "## 關鍵證據",
        ]
        body.extend(f"- {item}" for item in analysis.get("evidence", []))
        body.extend(["", "## 投影片大綱"])
        for slide in slide_outline:
            body.append(f"### {slide['title']}")
            body.extend(f"- {bullet}" for bullet in slide["bullets"])
        return "\n".join(frontmatter + body)

    def _build_slide_markdown(self, title: str, slide_outline: List[Dict]) -> str:
        lines = [f"# {title}", ""]
        for index, slide in enumerate(slide_outline, start=1):
            lines.append(f"## Slide {index}: {slide['title']}")
            lines.extend(f"- {bullet}" for bullet in slide["bullets"])
            lines.append("")
        return "\n".join(lines).strip()

    def _build_sources_slide(
        self,
        document: Dict,
        analysis: Dict,
        provider: str,
        model_name: str,
        target_language: str,
        user_prompt: str = None,
    ) -> Dict:
        bullets: List[str] = []
        title = document.get("title") or "Untitled document"
        bullets.append(f"原始文件 / Source title: {title}")
        if document.get("url"):
            bullets.append(f"來源網址 / URL: {document['url']}")
        if document.get("file_name"):
            bullets.append(f"原始檔名 / File name: {document['file_name']}")
        if document.get("source_name") or document.get("source_type"):
            bullets.append(
                f"來源類型 / Source: {document.get('source_type') or 'n/a'} via {document.get('source_name') or 'n/a'}"
            )
        if document.get("language"):
            bullets.append(f"原文語言 / Original language: {document['language']}")
        if document.get("country") or document.get("industry"):
            bullets.append(
                f"地區 / 產業: {document.get('country') or 'n/a'} · {document.get('industry') or 'n/a'}"
            )
        if document.get("published_date"):
            bullets.append(f"原始發布 / Published: {document['published_date']}")
        if document.get("created_at"):
            bullets.append(f"匯入時間 / Ingested at: {document['created_at']}")
        if analysis.get("risk_level"):
            risk_tags = ", ".join(analysis.get("risk_tags", []) or [])
            bullets.append(
                f"風險判斷 / Risk: {analysis['risk_level']}" + (f" ({risk_tags})" if risk_tags else "")
            )
        bullets.append(
            f"分析模型 / Generated by: provider={provider}, model={model_name}, target={target_language}"
        )
        bullets.append(f"報告時間 / Report timestamp: {datetime.now().isoformat(timespec='seconds')}")
        if user_prompt:
            bullets.append("使用者意圖 / Research intent: " + (user_prompt[:240] + "…" if len(user_prompt) > 240 else user_prompt))
        return {
            "title": "資料來源 / Sources & Citations",
            "bullets": bullets,
        }

    def _build_action_items(self, analysis: Dict) -> List[str]:
        if analysis["risk_level"] == "High":
            return ["立即檢視申報義務與生效日期", "安排稅務與法務共同評估", "建立後續追蹤提醒"]
        if analysis["risk_level"] == "Medium":
            return ["確認是否影響既有流程", "納入本月監控清單", "整理受影響部門"]
        return ["持續追蹤後續公告", "保留本次摘要供日後比對"]

    def _build_pptx_file(self, title: str, slide_outline: List[Dict]) -> str:
        if Presentation is None:
            raise ValueError("python-pptx is not installed. Please run pip install -r requirements.txt")

        presentation = Presentation()
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(5.625)

        presenter = "Tax Monitor"
        affiliation = "Automated Tax Risk Monitoring"
        theme_color = RGBColor(44, 62, 107)
        font_cn = "Taipei Sans TC Beta"
        font_en = "Times New Roman"

        self._add_template_title_slide(
            presentation,
            {
                "title": title,
                "presenter": presenter,
                "affiliation": affiliation,
            },
            theme_color,
            font_cn,
            font_en,
        )

        content_slides = [
            {
                "title": "Executive Snapshot",
                "bullets": [
                    f"{item.get('title', 'Untitled')}: {self._first_bullet(item)}"
                    for item in slide_outline[:4]
                ],
                "font_size": 20,
            }
        ] + slide_outline

        for page_number, slide in enumerate(content_slides, start=2):
            self._add_template_content_slide(
                presentation,
                slide,
                page_number,
                presenter,
                theme_color,
                font_cn,
                font_en,
            )

        report_dir = self._get_report_dir()
        file_name = f"{self._sanitize_file_name(title)}.pptx"
        file_path = report_dir / file_name
        presentation.save(str(file_path))
        return str(file_path)

    def _add_template_footer(self, slide, presentation, page_number: str, presenter: str, theme_color, font_en: str):
        footer_height = Inches(0.3)
        footer_y = presentation.slide_height - footer_height

        date_box = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(2.0), footer_height)
        date_paragraph = date_box.text_frame.paragraphs[0]
        date_paragraph.text = datetime.now().strftime("%Y/%m/%d")
        date_paragraph.font.name = font_en
        date_paragraph.font.size = Pt(12)

        page_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            (presentation.slide_width - Cm(4.06)) / 2,
            presentation.slide_height - Cm(0.6),
            Cm(4.06),
            Cm(0.6),
        )
        page_rect.fill.solid()
        page_rect.fill.fore_color.rgb = theme_color
        page_rect.line.fill.background()

        page_box = slide.shapes.add_textbox(
            (presentation.slide_width - Cm(4.06)) / 2,
            presentation.slide_height - Cm(0.6),
            Cm(4.06),
            Cm(0.6),
        )
        page_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        page_paragraph = page_box.text_frame.paragraphs[0]
        page_paragraph.text = page_number
        page_paragraph.font.name = font_en
        page_paragraph.font.size = Pt(12)
        page_paragraph.font.color.rgb = RGBColor(255, 255, 255)
        page_paragraph.alignment = PP_ALIGN.CENTER

        presenter_box = slide.shapes.add_textbox(Inches(7.5), footer_y, Inches(2.0), footer_height)
        presenter_paragraph = presenter_box.text_frame.paragraphs[0]
        presenter_paragraph.text = presenter
        presenter_paragraph.font.name = font_en
        presenter_paragraph.font.size = Pt(12)
        presenter_paragraph.alignment = PP_ALIGN.RIGHT

    def _add_template_title_slide(self, presentation, slide_data: Dict, theme_color, font_cn: str, font_en: str):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            Inches(1.5),
            Inches(8.6),
            Inches(1.2),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = theme_color
        rect.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = slide_data.get("title", "Tax Monitor Report")
        title_paragraph.font.name = font_cn
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
        title_paragraph.alignment = PP_ALIGN.CENTER

        presenter_box = slide.shapes.add_textbox(Inches(0), Inches(3.2), presentation.slide_width, Inches(0.5))
        presenter_paragraph = presenter_box.text_frame.paragraphs[0]
        presenter_paragraph.text = slide_data.get("presenter", "Tax Monitor")
        presenter_paragraph.font.name = font_cn
        presenter_paragraph.font.size = Pt(28)
        presenter_paragraph.alignment = PP_ALIGN.CENTER

        affiliation_box = slide.shapes.add_textbox(Inches(0), Inches(3.8), presentation.slide_width, Inches(0.4))
        affiliation_paragraph = affiliation_box.text_frame.paragraphs[0]
        affiliation_paragraph.text = slide_data.get("affiliation", "")
        affiliation_paragraph.font.name = font_cn
        affiliation_paragraph.font.size = Pt(22)
        affiliation_paragraph.font.italic = True
        affiliation_paragraph.font.color.rgb = theme_color
        affiliation_paragraph.alignment = PP_ALIGN.CENTER

        date_box = slide.shapes.add_textbox(Inches(0), Inches(4.4), presentation.slide_width, Inches(0.5))
        date_paragraph = date_box.text_frame.paragraphs[0]
        date_paragraph.text = datetime.now().strftime("%B %d, %Y")
        date_paragraph.font.name = font_en
        date_paragraph.font.size = Pt(20)
        date_paragraph.alignment = PP_ALIGN.CENTER

        self._add_template_footer(
            slide,
            presentation,
            "1",
            slide_data.get("presenter", "Tax Monitor"),
            theme_color,
            font_en,
        )

    def _add_template_content_slide(
        self,
        presentation,
        slide_data: Dict,
        page_number: int,
        presenter: str,
        theme_color,
        font_cn: str,
        font_en: str,
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        content_title_height = Cm(1.8)

        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            presentation.slide_width,
            content_title_height,
        )
        header.fill.solid()
        header.fill.fore_color.rgb = theme_color
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0), Inches(0), presentation.slide_width, content_title_height)
        title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_box.text_frame.margin_left = Inches(0.3)
        title_paragraph = title_box.text_frame.paragraphs[0]
        title_paragraph.text = str(slide_data.get("title", "Untitled"))
        title_paragraph.font.name = font_cn
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
        title_paragraph.alignment = PP_ALIGN.LEFT

        content_top = content_title_height + Inches(0.25)
        safe_height = presentation.slide_height - Inches(0.8) - content_top
        text_box = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(9.0), safe_height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()

        bullets = slide_data.get("bullets") or ["No detail generated."]
        font_size = int(slide_data.get("font_size", 22))
        for index, bullet in enumerate(bullets[:7]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = f"• {str(bullet)}"
            paragraph.level = 0
            paragraph.font.name = font_cn
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor(34, 44, 56)
            paragraph.space_after = Pt(12)

        self._add_template_footer(slide, presentation, str(page_number), presenter, theme_color, font_en)

    def _add_title_slide(self, presentation, title: str):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_background(slide, RGBColor(247, 249, 252))
        self._add_top_bar(slide, "Tax Monitor Risk Brief")

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(8.8), Inches(1.6))
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(24, 33, 43)

        subtitle = slide.shapes.add_textbox(Inches(0.72), Inches(3.05), Inches(7.6), Inches(0.45))
        subtitle_frame = subtitle.text_frame
        subtitle_frame.text = "Automated tax research, risk triage, and management-ready summary"
        subtitle_frame.paragraphs[0].font.size = Pt(15)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(85, 99, 115)

        badge = slide.shapes.add_shape(1, Inches(0.72), Inches(4.1), Inches(2.35), Inches(0.48))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(28, 108, 92)
        badge.line.color.rgb = RGBColor(28, 108, 92)
        badge.text_frame.text = "Generated by Tax Monitor"
        badge.text_frame.paragraphs[0].font.size = Pt(12)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_summary_slide(self, presentation, slide_outline: List[Dict]):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_background(slide, RGBColor(255, 255, 255))
        self._add_top_bar(slide, "Executive Snapshot")
        self._add_slide_title(slide, "Key Findings")

        left = Inches(0.7)
        top = Inches(1.55)
        width = Inches(11.9)
        height = Inches(4.9)
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        for index, item in enumerate(slide_outline[:5]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = f"{item.get('title', 'Untitled')}: {self._first_bullet(item)}"
            paragraph.font.size = Pt(17)
            paragraph.font.color.rgb = RGBColor(34, 44, 56)
            paragraph.space_after = Pt(12)

    def _add_content_slide(self, presentation, slide_data: Dict):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_background(slide, RGBColor(255, 255, 255))
        self._add_top_bar(slide, "Tax Monitor")
        self._add_slide_title(slide, slide_data.get("title", "Untitled"))

        box = slide.shapes.add_textbox(Inches(0.85), Inches(1.6), Inches(11.4), Inches(5.2))
        frame = box.text_frame
        frame.clear()
        bullets = slide_data.get("bullets") or ["No bullet points generated."]
        for index, bullet in enumerate(bullets[:6]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(34, 44, 56)
            paragraph.space_after = Pt(10)

    def _set_background(self, slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_top_bar(self, slide, label: str):
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(28, 108, 92)
        bar.line.color.rgb = RGBColor(28, 108, 92)
        text = slide.shapes.add_textbox(Inches(0.7), Inches(0.1), Inches(5.5), Inches(0.28))
        text.text_frame.text = label
        paragraph = text.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    def _add_slide_title(self, slide, title: str):
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(11.8), Inches(0.65))
        title_box.text_frame.text = title
        paragraph = title_box.text_frame.paragraphs[0]
        paragraph.font.size = Pt(25)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(24, 33, 43)

    def _first_bullet(self, slide_data: Dict) -> str:
        bullets = slide_data.get("bullets") or []
        return str(bullets[0]) if bullets else "No detail generated."

    def _get_report_dir(self) -> Path:
        for candidate in (REPORT_DIR, FALLBACK_REPORT_DIR, TEMP_REPORT_DIR):
            try:
                candidate.mkdir(parents=True, exist_ok=True)
                return candidate
            except PermissionError:
                continue
        raise PermissionError("No writable report directory is available")

    def _sanitize_file_name(self, title: str) -> str:
        cleaned = re.sub(r"[\\\\/:*?\"<>|]+", "_", title).strip()
        return cleaned[:80] or "tax_monitor_report"
